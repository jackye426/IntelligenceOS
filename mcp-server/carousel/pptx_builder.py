"""Build filled carousel PPTX decks from template definitions."""

from __future__ import annotations

import json
import re
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .text_fit import adjust_y_for_valign, fit_text_to_box

TEMPLATES_PATH = Path(__file__).resolve().parent / "templates.json"
ASSETS_DIR = Path(__file__).resolve().parent / "assets"
SLIDE_W_IN = 10.0
SLIDE_H_IN = 12.5

_WS_RE = re.compile(r"\s+")


@dataclass
class CarouselContent:
    topic: str
    template_id: str
    hook: str
    hook_subtitle: str | None
    slides: list[dict[str, str]]
    cta: str
    caption: str
    track_label: str | None = None


def _safe_filename(s: str) -> str:
    s = _WS_RE.sub(" ", (s or "").strip())
    s = re.sub(r"[^\w\- ]+", "", s, flags=re.UNICODE)
    s = s.replace(" ", "_")
    return (s or "carousel")[:100]


def load_templates() -> list[dict[str, Any]]:
    data = json.loads(TEMPLATES_PATH.read_text(encoding="utf-8"))
    return list(data.get("templates") or [])


def get_template(template_id: str) -> dict[str, Any]:
    for t in load_templates():
        if t.get("id") == template_id:
            return t
    known = [t.get("id") for t in load_templates()]
    raise ValueError(f"Unknown template_id={template_id!r}. Known: {known}")


def _align(align: str) -> PP_ALIGN:
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)


def _valign(valign: str) -> MSO_ANCHOR:
    return {
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)


def _font_name(template: dict[str, Any], font_key: str) -> str:
    fonts = template.get("fonts") or {}
    return str(fonts.get(font_key) or fonts.get("body") or "Work Sans")


def _set_bg(slide, template: dict[str, Any], bg_image: Path | None) -> None:
    bg = template.get("background") or {}
    if bg_image and bg_image.exists():
        sw, sh = slide.part.presentation.slide_width, slide.part.presentation.slide_height
        slide.shapes.add_picture(str(bg_image), 0, 0, width=sw, height=sh)
        return
    if bg.get("type") == "color":
        rgb = bg.get("rgb") or [255, 255, 255]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def _add_fitted_textbox(
    slide,
    *,
    zone: dict[str, Any],
    text: str,
    template: dict[str, Any],
) -> dict[str, Any]:
    role = str(zone.get("role") or "body")
    font_key = str(zone.get("font") or "body")
    x = float(zone.get("x", 1.0))
    y = float(zone.get("y", 1.0))
    w = float(zone.get("w", 8.0))
    h = float(zone.get("h", 2.0))
    max_pt = int(zone.get("max_pt", 24))
    min_pt = int(zone.get("min_pt", 12))
    color = zone.get("color") or [0, 0, 0]
    align = str(zone.get("align") or "left")
    valign = str(zone.get("valign") or "top")

    fit = fit_text_to_box(
        text,
        width_in=w,
        height_in=h,
        max_pt=max_pt,
        min_pt=min_pt,
        font_role=role if role in {"headline", "cta", "meta"} else "body",
    )
    y_adj = adjust_y_for_valign(
        y_in=y,
        box_h_in=h,
        line_count=fit.line_count,
        font_size_pt=fit.font_size_pt,
        font_role=role,
        valign=valign,
    )

    tb = slide.shapes.add_textbox(Inches(x), Inches(y_adj), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = _valign(valign)

    p = tf.paragraphs[0]
    p.text = fit.text
    p.alignment = _align(align)
    p.line_spacing = 1.1
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = _font_name(template, font_key)
    run.font.size = Pt(fit.font_size_pt)
    run.font.color.rgb = RGBColor(int(color[0]), int(color[1]), int(color[2]))
    if role in {"headline", "cta", "meta", "caption_label"}:
        run.font.bold = role in {"headline", "cta", "meta", "caption_label"}

    return {
        "role": role,
        "font_size_pt": fit.font_size_pt,
        "line_count": fit.line_count,
        "overflow": fit.overflow,
        "y_adjusted_in": round(y_adj, 2),
    }


def _render_layout_slide(
    prs: Presentation,
    *,
    template: dict[str, Any],
    layout_name: str,
    field_map: dict[str, str],
    bg_image: Path | None,
) -> list[dict[str, Any]]:
    layouts = template.get("layouts") or {}
    layout = layouts.get(layout_name)
    if not layout:
        raise ValueError(f"Template {template.get('id')} missing layout {layout_name}")

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, template, bg_image)

    metrics: list[dict[str, Any]] = []
    for zone in layout.get("zones") or []:
        role = str(zone.get("role") or "")
        text = field_map.get(role, "")
        if not text and role != "caption_label":
            continue
        if role == "caption_label":
            text = text or "Caption"
        m = _add_fitted_textbox(slide, zone=zone, text=text, template=template)
        metrics.append(m)
    return metrics


def build_carousel_pptx(
    content: CarouselContent,
    *,
    out_dir: Path,
    bg_image: Path | None = None,
) -> tuple[Path, dict[str, Any]]:
    template = get_template(content.template_id)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Optional per-template background image: assets/<template_id>.jpg
    template_bg = ASSETS_DIR / "backgrounds" / f"{content.template_id}.jpg"
    effective_bg = bg_image if bg_image and bg_image.exists() else (
        template_bg if template_bg.exists() else None
    )

    all_metrics: dict[str, Any] = {"slides": []}

    # Title / planning slide
    title_fields = {
        "meta": f"### {content.topic}",
        "caption_label": "Caption",
        "caption": content.caption,
    }
    if content.track_label:
        title_fields["meta"] = f"### {content.topic}\n{content.track_label}"
    title_metrics = _render_layout_slide(
        prs,
        template=template,
        layout_name="title",
        field_map=title_fields,
        bg_image=None,  # title slide uses solid bg for readability
    )
    all_metrics["slides"].append({"type": "title", "text_metrics": title_metrics})

    # Hook
    hook_fields = {
        "headline": content.hook,
        "subtitle": content.hook_subtitle or "",
    }
    hook_metrics = _render_layout_slide(
        prs,
        template=template,
        layout_name="hook",
        field_map=hook_fields,
        bg_image=effective_bg,
    )
    all_metrics["slides"].append({"type": "hook", "text_metrics": hook_metrics})

    # Body slides
    for i, sl in enumerate(content.slides, start=1):
        body_fields = {
            "headline": sl.get("main_text", ""),
            "body": sl.get("subtext", ""),
        }
        body_metrics = _render_layout_slide(
            prs,
            template=template,
            layout_name="body",
            field_map=body_fields,
            bg_image=effective_bg,
        )
        all_metrics["slides"].append({"type": "body", "index": i, "text_metrics": body_metrics})

    # CTA
    cta_metrics = _render_layout_slide(
        prs,
        template=template,
        layout_name="cta",
        field_map={"cta": content.cta},
        bg_image=effective_bg,
    )
    all_metrics["slides"].append({"type": "cta", "text_metrics": cta_metrics})

    out_dir.mkdir(parents=True, exist_ok=True)
    fname = f"{_safe_filename(content.topic)}_{content.template_id}_{uuid.uuid4().hex[:8]}.pptx"
    out_path = out_dir / fname
    prs.save(str(out_path))

    overflow_any = any(
        m.get("overflow")
        for slide in all_metrics["slides"]
        for m in slide.get("text_metrics", [])
    )
    all_metrics["overflow_warning"] = overflow_any
    all_metrics["template_id"] = content.template_id
    all_metrics["slide_count"] = len(content.slides) + 3  # title + hook + cta

    return out_path, all_metrics
